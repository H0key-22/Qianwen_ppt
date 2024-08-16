import random
import json
from pptx import Presentation
from dashscope import Generation
import dashscope

# 调用阿里通义千问大模型API
def qianwen_llm(query, history=[], user_stop_words=[]):

    # 读取个人api
    dashscope.api_key = "API-KEY"

    # 构造LLM输入
    messages = [{'role': 'system', 'content': 'You are a helpful assistant.'}]
    for hist in history:
        messages.append({'role': 'user', 'content': hist[0]})
        messages.append({'role': 'assistant', 'content': hist[1]})
    messages.append({'role': 'user', 'content': query})

    # 调用API
    response = Generation.call(model="qwen-turbo",
                               messages=messages,
                               # 设置随机数种子seed，如果没有设置，则随机数种子默认为1234
                               seed=random.randint(1, 10000),
                               temperature=0.8,
                               top_p=0.8,
                               top_k=50,
                               # 将输出设置为"message"格式
                               result_format='message')
    print(response)
    return response['output']['choices'][0]['message']['content']


# 生成PPT内容
def generate_ppt_content(topic, pages):
    # 输出格式
    output_format = json.dumps({
        "title": "example title",
        "pages": [
            {
                "title": "title for page 1",
                "content": [
                    {
                        "title": "title for paragraph 1",
                        "description": "detail for paragraph 1",
                    },
                    {
                        "title": "title for paragraph 2",
                        "description": "detail for paragraph 2",
                    },
                ],
            },
            {
                "title": "title for page 2",
                "content": [
                    {
                        "title": "title for paragraph 1",
                        "description": "detail for paragraph 1",
                    },
                    {
                        "title": "title for paragraph 2",
                        "description": "detail for paragraph 2",
                    },
                    {
                        "title": "title for paragraph 3",
                        "description": "detail for paragraph 3",
                    },
                ],
            },
        ],
    }, ensure_ascii=True)

    # prompt
    prompt = f'''我要准备1个关于{topic}的PPT，要求一共写{pages}页，请你根据主题生成详细内容，不要省略。
    按这个JSON格式输出{output_format}，只能返回JSON，且JSON不要用```包裹，内容要用中文。'''

    # print(prompt)

    # 调用llm生成PPT内容
    ppt_content = json.loads(qianwen_llm(prompt))
    return ppt_content


# 生成PPT文件
def generate_ppt_file(topic, ppt_content):
    ppt = Presentation()

    # PPT首页
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])  # title&subtitle layout
    slide.placeholders[0].text = ppt_content['title']
    slide.placeholders[1].text = "通义千问72B"

    # 内容页
    print('总共%d页...' % len(ppt_content['pages']))
    for i, page in enumerate(ppt_content['pages']):
        print('生成第%d页:%s' % (i + 1, page['title']))
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # title&content layout
        # 标题
        slide.placeholders[0].text = page['title']
        # 正文
        for sub_content in page['content']:
            print(sub_content)
            # 一级正文
            sub_title = slide.placeholders[1].text_frame.add_paragraph()
            sub_title.text, sub_title.level = sub_content['title'], 1
            # 二级正文
            sub_description = slide.placeholders[1].text_frame.add_paragraph()
            sub_description.text, sub_description.level = sub_content['description'], 2

    ppt.save('%s.pptx' % topic)


if __name__ == '__main__':
    while True:
        # 输入需求
        topic = input('输入主题:')
        pages = int(input('输入页数:'))
        # 生成PPT内容
        ppt_content = generate_ppt_content(topic, pages)
        # 生成PPT文件
        generate_ppt_file(topic, ppt_content)